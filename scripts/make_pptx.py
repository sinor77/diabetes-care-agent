"""Generate the MMU Showcase PowerPoint with a clean light theme."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Clean light theme palette (matching MMU prep guide style)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF8, 0xFA, 0xFC)
TEXT = RGBColor(0x0F, 0x17, 0x2A)        # near-black
MUTED = RGBColor(0x64, 0x74, 0x8B)        # gray
ACCENT = RGBColor(0x16, 0xA3, 0x4A)       # green primary
ACCENT_LIGHT = RGBColor(0xDC, 0xFC, 0xE7) # light green
ACCENT_BG = RGBColor(0xAC, 0xFF, 0xCB)    # mint logo bg
BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

def add_blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    return slide

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font_name
    return tb

def add_brand_header(slide, segment_text=""):
    # Logo
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.3), Inches(0.5), Inches(0.5))
    logo.fill.solid(); logo.fill.fore_color.rgb = ACCENT_BG
    logo.line.fill.background()
    logo.adjustments[0] = 0.3
    add_text(slide, Inches(0.4), Inches(0.3), Inches(0.5), Inches(0.5), "🩺", size=18, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Brand text
    add_text(slide, Inches(1.0), Inches(0.32), Inches(6), Inches(0.5),
             "Diabetes Control AI", size=12, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    if segment_text:
        add_text(slide, Inches(8), Inches(0.32), Inches(5), Inches(0.5), segment_text, size=10, color=MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def add_slide_num(slide, num, total=12):
    add_text(slide, Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3), f"{num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def add_card(slide, x, y, w, h, fill=CARD_BG, border=BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    card.adjustments[0] = 0.04
    return card

def add_arch_box(slide, x, y, w, h, text, color, accent):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)
    box.adjustments[0] = 0.1
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05); tf.margin_top = tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"

def add_arrow(slide, x, y, w=Inches(0.35), h=Inches(0.25)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()

# ====================== SLIDE 1: COVER ======================
s = add_blank_slide()
# Top accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

# Big logo box
logo_big = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.67), Inches(1.5), Inches(2), Inches(2))
logo_big.fill.solid(); logo_big.fill.fore_color.rgb = ACCENT_BG; logo_big.line.fill.background()
logo_big.adjustments[0] = 0.2
add_text(s, Inches(5.67), Inches(1.5), Inches(2), Inches(2), "🩺", size=80, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.0), "Diabetes Control AI", size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.5), "Your personal endocrinologist, available 24/7", size=20, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.4), "Digital Health  ·  Powered by Amazon Bedrock (Claude 3 Haiku)", size=14, color=TEXT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
         "AWS CendekiAwan 2.0  ·  MMU Showcase  ·  24 June 2026", size=11, color=MUTED, align=PP_ALIGN.CENTER)
add_slide_num(s, 1)

# ====================== SLIDE 2: ARCHITECTURE ======================
s = add_blank_slide()
add_brand_header(s, "Architecture")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "How It All Connects", size=32, bold=True, color=TEXT)
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4), "14 AWS Services  ·  Fully Serverless  ·  Zero Server Management", size=14, color=MUTED)

# Architecture canvas
canvas = add_card(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.85))

# Row 1: Users
add_arch_box(s, Inches(2.5), Inches(2.4), Inches(2), Inches(0.45), "👤 Patient", BLUE, BLUE)
add_arch_box(s, Inches(5.7), Inches(2.4), Inches(2), Inches(0.45), "👨‍⚕️ Doctor", PURPLE, PURPLE)
add_arch_box(s, Inches(8.9), Inches(2.4), Inches(2), Inches(0.45), "🔐 Admin", RED, RED)

# Row 2: CDN + S3
add_arch_box(s, Inches(3.5), Inches(3.1), Inches(2.6), Inches(0.45), "Amazon CloudFront (HTTPS CDN)", RGBColor(0x0E, 0xA5, 0xE9), RGBColor(0x0E, 0xA5, 0xE9))
add_arrow(s, Inches(6.2), Inches(3.18))
add_arch_box(s, Inches(6.7), Inches(3.1), Inches(2.6), Inches(0.45), "Amazon S3 (Static Site)", RGBColor(0x0E, 0xA5, 0xE9), RGBColor(0x0E, 0xA5, 0xE9))

# Row 3: Auth + API + Lambda
add_arch_box(s, Inches(2.0), Inches(3.85), Inches(2.4), Inches(0.45), "Cognito (Auth & Roles)", PURPLE, PURPLE)
add_arrow(s, Inches(4.5), Inches(3.93))
add_arch_box(s, Inches(5.05), Inches(3.85), Inches(2.4), Inches(0.45), "API Gateway (REST)", AMBER, AMBER)
add_arrow(s, Inches(7.55), Inches(3.93))
add_arch_box(s, Inches(8.1), Inches(3.85), Inches(2.4), Inches(0.45), "AWS Lambda (Python)", ACCENT, ACCENT)

# Row 4: AI + Data services
add_arch_box(s, Inches(0.8), Inches(4.6), Inches(2.4), Inches(0.6), "Amazon Bedrock\nClaude 3 Haiku", RGBColor(0xEA, 0x58, 0x0C), RGBColor(0xEA, 0x58, 0x0C))
add_arch_box(s, Inches(3.4), Inches(4.6), Inches(2.4), Inches(0.6), "Amazon DynamoDB\n3 tables", BLUE, BLUE)
add_arch_box(s, Inches(6.0), Inches(4.6), Inches(2.4), Inches(0.6), "Amazon Textract\n(Lab OCR)", PURPLE, PURPLE)
add_arch_box(s, Inches(8.6), Inches(4.6), Inches(2.4), Inches(0.6), "Amazon Polly\n(Text-to-Speech)", RGBColor(0xDB, 0x27, 0x77), RGBColor(0xDB, 0x27, 0x77))
add_arch_box(s, Inches(11.2), Inches(4.6), Inches(1.5), Inches(0.6), "Step\nFunctions", AMBER, AMBER)

# Row 5: Supporting
add_arch_box(s, Inches(0.8), Inches(5.4), Inches(2.4), Inches(0.5), "Amazon SES (Email)", ACCENT, ACCENT)
add_arch_box(s, Inches(3.4), Inches(5.4), Inches(2.4), Inches(0.5), "EventBridge (Alerts)", PURPLE, PURPLE)
add_arch_box(s, Inches(6.0), Inches(5.4), Inches(2.4), Inches(0.5), "CloudWatch (Monitoring)", RED, RED)
add_arch_box(s, Inches(8.6), Inches(5.4), Inches(4.1), Inches(0.5), "Comprehend Medical (Healthcare NLP)", RGBColor(0x06, 0xB6, 0xD4), RGBColor(0x06, 0xB6, 0xD4))

add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.6),
         "Flow:  User → CloudFront → API Gateway → Lambda → Bedrock Agent (Claude 3 Haiku) → response  +  persisted to DynamoDB",
         size=11, color=TEXT, bold=True)

add_slide_num(s, 2)

# ====================== SLIDE 3: PROBLEM ======================
s = add_blank_slide()
add_brand_header(s, "Problem  ·  0:00 – 1:00")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "Meet Ahmed", size=32, bold=True, color=TEXT)

# Persona card
persona = add_card(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.6))
# Left accent bar on the card
accent_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(0.08), Inches(2.6))
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = ACCENT; accent_bar.line.fill.background()

add_text(s, Inches(0.85), Inches(1.85), Inches(11.5), Inches(0.45), "Ahmed Hassan, 52  ·  Bus driver in Klang Valley", size=20, bold=True, color=TEXT)
add_text(s, Inches(0.85), Inches(2.35), Inches(11.5), Inches(0.6),
         "Type 2 diabetes for 8 years.  HbA1c stuck at 7.8% (target: 6.5%).",
         size=14, color=TEXT)
add_text(s, Inches(0.85), Inches(2.78), Inches(11.5), Inches(0.6),
         "Next endocrinologist appointment is 14 weeks away at HKL.",
         size=14, color=TEXT)
add_text(s, Inches(0.85), Inches(3.25), Inches(11.5), Inches(0.6),
         "Last week his sugar hit 245 mg/dL after dinner. He doesn't know if that's dangerous, what to eat tomorrow, or whether to take extra metformin.",
         size=12, color=MUTED)
add_text(s, Inches(0.85), Inches(3.85), Inches(11.5), Inches(0.4),
         "\"I don't want to call the hospital for every question — but I don't want to die either.\"", size=12, color=MUTED, bold=False)

# Stats row
def add_stat(slide, x, big, label):
    add_card(slide, x, Inches(4.7), Inches(3.95), Inches(2.0))
    add_text(slide, x, Inches(4.85), Inches(3.95), Inches(0.85), big, size=42, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), Inches(5.85), Inches(3.55), Inches(0.7), label, size=11, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

add_stat(s, Inches(0.5), "3.9M", "Malaysians with diabetes (1 in 5 adults)")
add_stat(s, Inches(4.7), "120+", "days avg wait for endocrinology in public hospitals")
add_stat(s, Inches(8.85), "73%", "don't reach HbA1c target — preventable complications")

add_slide_num(s, 3)

# ====================== SLIDE 4: SOLUTION OVERVIEW ======================
s = add_blank_slide()
add_brand_header(s, "Solution  ·  1:00 – 2:00")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "An AI Endocrinologist in Your Pocket", size=30, bold=True, color=TEXT)
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4), "Three connected experiences, one intelligent agent", size=14, color=MUTED)

# Three solution cards
def solution_card(slide, x, icon, title, desc):
    card = add_card(slide, x, Inches(2.1), Inches(4.05), Inches(1.9), fill=ACCENT_LIGHT, border=ACCENT)
    add_text(slide, x + Inches(0.3), Inches(2.25), Inches(3.7), Inches(0.5), f"{icon}  {title}", size=15, bold=True, color=TEXT)
    add_text(slide, x + Inches(0.3), Inches(2.85), Inches(3.7), Inches(1.05), desc, size=11, color=MUTED)

solution_card(s, Inches(0.5), "👤", "Patient App", "Profile, meal logger, lab analyzer, AI chat coach. Get instant insights when symptoms appear.")
solution_card(s, Inches(4.65), "👨‍⚕️", "Doctor Dashboard", "Verified physicians review referred patients with charts, lab images, and AI-flagged risks.")
solution_card(s, Inches(8.8), "🔐", "Admin Panel", "Verify doctor credentials before they appear to patients — clinical safety layer.")

# Why AI card (warning style)
add_card(s, Inches(0.5), Inches(4.25), Inches(12.3), Inches(2.85), fill=RGBColor(0xFE, 0xF3, 0xC7), border=AMBER)
add_text(s, Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.45), "⚠️  Why this needs AI — not a form or database", size=15, bold=True, color=AMBER)

ai_reasons = [
    ("📷 Vision understanding", "Read a Malaysian Pathology Lab photo — Claude extracts HbA1c, lipids, kidney markers automatically"),
    ("🧠 Contextual reasoning", "\"I had nasi lemak this morning, sugar is 220\" → AI knows portion sizes, glycemic load, what to do NOW"),
    ("💬 Natural conversation", "Ahmed types in plain English (or BM) — no medical jargon, no rigid form fields"),
    ("🔍 Cross-references unstructured data", "Profile + labs + chat history = personalized clinical guidance no form can deliver"),
]
y = 4.95
for icon_title, desc in ai_reasons:
    add_text(s, Inches(0.85), Inches(y), Inches(3.5), Inches(0.4), icon_title, size=12, bold=True, color=TEXT)
    add_text(s, Inches(4.4), Inches(y), Inches(8.2), Inches(0.4), desc, size=11, color=MUTED)
    y += 0.5

add_slide_num(s, 4)

# ====================== SLIDE 5: TECH STACK ======================
s = add_blank_slide()
add_brand_header(s, "Solution  ·  AWS Architecture")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "14 AWS Services Working Together", size=30, bold=True, color=TEXT)

# Table-like layout
rows = [
    ("Layer", "Service", "Purpose"),
    ("AI Brain", "Amazon Bedrock", "Claude 3 Haiku — $0.25/1M tokens, low latency, multimodal vision"),
    ("Compute", "Lambda · API Gateway", "Serverless API with 14+ endpoints, auto-scales"),
    ("Auth", "Amazon Cognito", "Role-based: Patient / Expert groups, email verification"),
    ("Data", "DynamoDB · S3", "3 tables: profiles, health-logs, referrals; S3 hosts frontend"),
    ("Specialized AI", "Textract · Comprehend Medical · Polly", "Lab OCR · medical NLP · neural text-to-speech"),
    ("Workflow", "Step Functions · EventBridge", "Multi-step analysis pipelines · threshold alerts"),
    ("Delivery", "CloudFront · SES", "HTTPS CDN · email reports"),
    ("Operations", "CloudWatch", "Live dashboard with API metrics"),
]
y = 1.7
for i, (layer, svc, purpose) in enumerate(rows):
    is_header = (i == 0)
    if is_header:
        add_text(s, Inches(0.5), Inches(y), Inches(2.5), Inches(0.4), layer, size=10, bold=True, color=MUTED)
        add_text(s, Inches(3.0), Inches(y), Inches(3.5), Inches(0.4), svc, size=10, bold=True, color=MUTED)
        add_text(s, Inches(6.5), Inches(y), Inches(6.3), Inches(0.4), purpose, size=10, bold=True, color=MUTED)
    else:
        add_text(s, Inches(0.5), Inches(y), Inches(2.5), Inches(0.4), layer, size=12, bold=True, color=ACCENT)
        add_text(s, Inches(3.0), Inches(y), Inches(3.5), Inches(0.4), svc, size=12, bold=True, color=TEXT)
        add_text(s, Inches(6.5), Inches(y), Inches(6.3), Inches(0.4), purpose, size=11, color=MUTED)
    # Underline
    line = s.shapes.add_connector(1, Inches(0.5), Inches(y + 0.45), Inches(12.8), Inches(y + 0.45))
    line.line.color.rgb = BORDER; line.line.width = Pt(0.5)
    y += 0.55

add_text(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
         "Flow:  Patient input → API Gateway → Lambda → Bedrock Agent → response + DynamoDB → Doctor sees real-time",
         size=12, bold=True, color=ACCENT)
add_slide_num(s, 5)

# ====================== SLIDE 6: LIVE DEMO INTRO ======================
s = add_blank_slide()
add_brand_header(s, "Live Demo  ·  2:00 – 5:00")
add_text(s, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.7), "Try It Yourself", size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.4), "Open the link  ·  Sign up takes 30 seconds", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# URL box
url_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.7))
url_box.fill.solid(); url_box.fill.fore_color.rgb = ACCENT_LIGHT; url_box.line.color.rgb = ACCENT; url_box.line.width = Pt(1.5)
url_box.adjustments[0] = 0.15
add_text(s, Inches(2.5), Inches(2.5), Inches(8.3), Inches(0.7), "https://d3onijdn12lthk.cloudfront.net", size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name="Consolas")

# Demo accounts cards
def account_card(slide, x, color_role, title, email, password):
    card = add_card(slide, x, Inches(3.7), Inches(4.05), Inches(2.6))
    add_text(slide, x + Inches(0.25), Inches(3.85), Inches(3.6), Inches(0.4), title, size=14, bold=True, color=color_role)
    add_text(slide, x + Inches(0.25), Inches(4.4), Inches(3.6), Inches(0.4), "Email:", size=10, color=MUTED)
    add_text(slide, x + Inches(0.25), Inches(4.7), Inches(3.6), Inches(0.4), email, size=11, bold=True, color=TEXT, font_name="Consolas")
    add_text(slide, x + Inches(0.25), Inches(5.2), Inches(3.6), Inches(0.4), "Password:", size=10, color=MUTED)
    add_text(slide, x + Inches(0.25), Inches(5.5), Inches(3.6), Inches(0.4), password, size=11, bold=True, color=TEXT, font_name="Consolas")

account_card(s, Inches(0.5), ACCENT, "Demo Patient", "patient.test@diabetes.demo", "Patient123!")
account_card(s, Inches(4.65), BLUE, "Demo Doctor (verified)", "doctor.test@diabetes.demo", "Doctor123!")
account_card(s, Inches(8.8), PURPLE, "Admin Panel  /admin.html", "Password only", "diabetes2025admin")

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
         "Pre-loaded with realistic data: 7-day HbA1c history, lab image, AI insights, doctor referral.",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
add_slide_num(s, 6)

# ====================== SLIDE 7: DEMO FLOW ======================
s = add_blank_slide()
add_brand_header(s, "Live Demo  ·  The Flow")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "The 3-Minute Demo Path", size=30, bold=True, color=TEXT)

steps = [
    ("1", "Sign in as Ahmed", "Pre-loaded profile loads from DynamoDB"),
    ("2", "Upload lab photo", "Claude vision extracts real values"),
    ("3", "Ask AI Coach", "\"What should I eat tonight?\""),
    ("4", "Send to Doctor", "Selective sharing checkboxes"),
    ("5", "Switch to Doctor", "See charts · gauges · AI insights"),
]
x = 0.5
for num, title, desc in steps:
    card = add_card(s, Inches(x), Inches(2.0), Inches(2.45), Inches(2.3))
    # Top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(2.45), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.0), Inches(2.25), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    add_text(s, Inches(x + 1.0), Inches(2.25), Inches(0.45), Inches(0.45), num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(x + 0.15), Inches(2.85), Inches(2.15), Inches(0.5), title, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(x + 0.15), Inches(3.45), Inches(2.15), Inches(0.7), desc, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    x += 2.55

# What you'll see live
add_card(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.4), fill=RGBColor(0xFE, 0xF3, 0xC7), border=AMBER)
add_text(s, Inches(0.75), Inches(4.85), Inches(11.8), Inches(0.45), "⚠️  What you'll see live", size=15, bold=True, color=AMBER)
live_items = [
    "📊  Real time-series charts from Ahmed's actual submission history",
    "🧠  AI-flagged anomalies — \"HbA1c trending down but LDL rising, recommend statin discussion\"",
    "📷  Lab photo Ahmed uploaded — full size, click to zoom",
    "🚨  EventBridge alert: \"Patient HbA1c above 8% for 3 weeks\" auto-triggered",
]
y = 5.4
for item in live_items:
    add_text(s, Inches(0.85), Inches(y), Inches(11.6), Inches(0.4), item, size=11, color=TEXT)
    y += 0.4

add_slide_num(s, 7)

# ====================== SLIDE 8: HARD QUESTIONS ======================
s = add_blank_slide()
add_brand_header(s, "Live Demo  ·  Hard Questions")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "How We Handle the Tough Questions", size=28, bold=True, color=TEXT)

def qa_card(slide, x, y, w, h, q, lines):
    add_card(slide, x, y, w, h)
    add_text(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.4), Inches(0.45), q, size=12, bold=True, color=AMBER)
    yy = y + Inches(0.65)
    for line in lines:
        add_text(slide, x + Inches(0.4), yy, w - Inches(0.55), Inches(0.32), "• " + line, size=10, color=TEXT)
        yy += Inches(0.32)

qa_card(s, Inches(0.5), Inches(1.7), Inches(6.15), Inches(2.6),
        "\"Which Bedrock model and why?\"",
        ["Claude 3 Haiku via APAC inference profile",
         "$0.25 per 1M input tokens — sustainable at scale",
         "Sub-second latency for chat responses",
         "Multimodal — reads lab images natively",
         "Available in ap-southeast-1 (data sovereignty)"])

qa_card(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(2.6),
        "\"How do you prevent hallucination?\"",
        ["Strict system prompt grounded in ADA Standards of Care",
         "Temperature 0.3 for clinical responses",
         "Cross-references DynamoDB patient facts",
         "Discrepancy detection: profile vs lab values flagged",
         "Hard guardrail: never adjust medication dosages"])

qa_card(s, Inches(0.5), Inches(4.45), Inches(6.15), Inches(2.6),
        "\"What about harmful inputs?\"",
        ["Bedrock Guardrails block self-harm content",
         "Crisis: redirects to 988 (suicide lifeline)",
         "DKA symptoms → \"Call 999 immediately\"",
         "System prompt forbids medication changes"])

qa_card(s, Inches(6.85), Inches(4.45), Inches(5.95), Inches(2.6),
        "\"What if input is unusual?\"",
        ["Tested 30+ edge cases (BM mixed, slang, missing fields)",
         "Fallback: AI says \"I'm not sure, please share with doctor\"",
         "Admin can audit chat logs in DynamoDB",
         "Lab vision falls back to text if image unclear"])

add_slide_num(s, 8)

# ====================== SLIDE 9: REAL WORLD IMPACT ======================
s = add_blank_slide()
add_brand_header(s, "Impact  ·  5:00 – 6:00")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "Built for Malaysia. Ready to Pilot.", size=30, bold=True, color=TEXT)

# Pilot partner card
add_card(s, Inches(0.5), Inches(1.7), Inches(6.15), Inches(2.5), fill=ACCENT_LIGHT, border=ACCENT)
add_text(s, Inches(0.7), Inches(1.85), Inches(5.8), Inches(0.4), "🏥  Pilot Partner", size=13, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(2.3), Inches(5.8), Inches(0.5), "Klinik Kesihatan Bandar Botanic", size=18, bold=True, color=TEXT)
add_text(s, Inches(0.7), Inches(2.85), Inches(5.8), Inches(1.0),
         "Klang Valley primary care clinic — 1,200 diabetes patients on file. Their endocrinologist sees patients every 3-4 months. Our app fills the in-between gap.",
         size=11, color=MUTED)
add_text(s, Inches(0.7), Inches(3.85), Inches(5.8), Inches(0.4),
         "Backup: MOH Klinik Komuniti, NADI, Persatuan Diabetes Malaysia",
         size=10, color=MUTED)

# Cost card
add_card(s, Inches(6.85), Inches(1.7), Inches(5.95), Inches(2.5))
add_text(s, Inches(7.05), Inches(1.85), Inches(5.6), Inches(0.4), "💰  Cost per User per Month", size=13, bold=True, color=ACCENT)
cost_rows = [
    ("Bedrock (200 chats × 1K tokens)", "$0.10"),
    ("Lambda (free tier)", "$0.00"),
    ("DynamoDB (free tier)", "$0.00"),
    ("SES (1 email)", "$0.0001"),
    ("CloudFront / S3", "$0.02"),
]
y = 2.3
for label, val in cost_rows:
    add_text(s, Inches(7.05), Inches(y), Inches(4.0), Inches(0.3), label, size=10, color=MUTED)
    add_text(s, Inches(11.0), Inches(y), Inches(1.6), Inches(0.3), val, size=10, color=TEXT, align=PP_ALIGN.RIGHT)
    y += 0.3
# Divider
divider = s.shapes.add_connector(1, Inches(7.05), Inches(3.85), Inches(12.6), Inches(3.85))
divider.line.color.rgb = BORDER; divider.line.width = Pt(0.75)
add_text(s, Inches(7.05), Inches(3.95), Inches(4.0), Inches(0.3), "Total per user / month", size=11, bold=True, color=TEXT)
add_text(s, Inches(11.0), Inches(3.95), Inches(1.6), Inches(0.3), "~$0.13", size=14, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# Three challenge cards (bottom)
def challenge_card(slide, x, title, body):
    add_card(slide, x, Inches(4.4), Inches(4.05), Inches(2.7))
    add_text(slide, x + Inches(0.25), Inches(4.55), Inches(3.6), Inches(0.5), title, size=12, bold=True, color=AMBER)
    add_text(slide, x + Inches(0.25), Inches(5.05), Inches(3.6), Inches(2.0), body, size=10, color=MUTED)

challenge_card(s, Inches(0.5), "⚠️  What breaks at 100K users?",
              "Bedrock TPM rate limits hit first. Solved with regional inference profiles + Lambda response caching for frequent queries (e.g., \"is 180 mg/dL high?\").")
challenge_card(s, Inches(4.65), "⚠️  Who gets left out?",
              "Elderly without smartphones. Plan: USSD/SMS interface via Amazon SNS for basic glucose logging — works on any feature phone.")
challenge_card(s, Inches(8.8), "⚠️  Why won't a big company copy?",
              "Local moat: BM/Manglish, Malaysian lab formats, halal meal database, integration with KKM clinic referral systems.")

add_slide_num(s, 9)

# ====================== SLIDE 10: TEAM ======================
s = add_blank_slide()
add_brand_header(s, "Team")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "Who Built What", size=30, bold=True, color=TEXT)

def team_card(slide, x, role, body):
    card = add_card(slide, x, Inches(2.0), Inches(4.05), Inches(3.5))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(4.05), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text(slide, x + Inches(0.3), Inches(2.3), Inches(3.6), Inches(0.5), role, size=15, bold=True, color=TEXT)
    add_text(slide, x + Inches(0.3), Inches(3.0), Inches(3.6), Inches(2.4), body, size=11, color=MUTED)

team_card(s, Inches(0.5), "Backend & Architecture",
          "Bedrock Agent setup with system prompts; Lambda handlers for 14 API endpoints; DynamoDB schema for profiles, health logs, referrals; IAM roles per service; Step Functions analysis pipeline; CloudFormation deployment.")
team_card(s, Inches(4.65), "Frontend & UX",
          "Three SPAs (patient, doctor, admin) sharing one backend; Tailwind styling; Chart.js visualizations (gauges, lines, bar, radar); dark/light theme; accessibility-first forms; persistent localStorage + DynamoDB sync.")
team_card(s, Inches(8.8), "AI Engineering & Clinical",
          "System prompts grounded in ADA Standards of Care; prompt engineering for clinical accuracy; lab vision integration via Bedrock Converse API; Comprehend Medical entity extraction; chart insight generation; safety guardrails.")

add_slide_num(s, 10)

# ====================== SLIDE 11: ROADMAP ======================
s = add_blank_slide()
add_brand_header(s, "What's Next")
add_text(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5), "If We Had 3 More Weeks", size=30, bold=True, color=TEXT)

roadmap = [
    ("🇲🇾  Bahasa Melayu support", "Amazon Translate so patients can chat in their first language"),
    ("📡  CGM integration", "Live glucose data from FreeStyle Libre / Dexcom via FHIR API"),
    ("📧  SES production access", "Emails go straight to inbox without verification step"),
    ("🌐  Custom domain", "diabetescare.my with proper DKIM/SPF for email trust"),
    ("📚  Bedrock Knowledge Base", "RAG against ADA Standards PDF for citation-backed answers"),
    ("📱  Native mobile app", "React Native wrapper, push notifications via SNS"),
    ("🤝  Hardware partnership", "Pilot real CGM devices with Abbott Malaysia"),
]
y = 1.85
for icon_title, desc in roadmap:
    # Bullet circle
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
    add_text(s, Inches(1.05), Inches(y), Inches(4.5), Inches(0.45), icon_title, size=13, bold=True, color=TEXT)
    add_text(s, Inches(5.6), Inches(y), Inches(7.2), Inches(0.45), desc, size=11, color=MUTED)
    y += 0.65

add_slide_num(s, 11)

# ====================== SLIDE 12: CLOSING ======================
s = add_blank_slide()
# Top accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.15))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.6), "One Sentence", size=22, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.0),
         "Diabetes Control AI helps Malaysians manage diabetes daily —\nwithout waiting 3 months for an endocrinologist appointment.",
         size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Stats line
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
         "Built on AWS in 6 weeks  ·  $0.13 per user per month  ·  Live now  ·  Ready to pilot",
         size=14, color=TEXT, align=PP_ALIGN.CENTER)

# URL box
url_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.7))
url_box.fill.solid(); url_box.fill.fore_color.rgb = ACCENT_LIGHT; url_box.line.color.rgb = ACCENT; url_box.line.width = Pt(1.5)
url_box.adjustments[0] = 0.15
add_text(s, Inches(3.5), Inches(5.3), Inches(6.3), Inches(0.7), "github.com/sinor77/diabetes-care-agent", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font_name="Consolas")

add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4), "Thank you  ·  Terima kasih  ·  شكراً", size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_slide_num(s, 12)

# Save
out = "presentation/DiabetesControl_AI_MMU_Showcase.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
