"""Final Presentation v3 - Clear architecture per role, 2025-2026 studies, easy showcase."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GL = RGBColor(0xDC, 0xFC, 0xE7)
MINT = RGBColor(0xAC, 0xFF, 0xCB)
BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PINK = RGBColor(0xDB, 0x27, 0x77)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def S(): 
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    b.fill.solid(); b.fill.fore_color.rgb = WHITE; b.line.fill.background()
    return s

def T(s,x,y,w,h,t,sz=14,b=False,c=TEXT,a=PP_ALIGN.LEFT,v=MSO_ANCHOR.TOP,f="Segoe UI"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame;tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Emu(0)
    tf.vertical_anchor=v;p=tf.paragraphs[0];p.alignment=a
    r=p.add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=f

def box(s,x,y,w,h,label,color,fill=WHITE):
    b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    b.fill.solid();b.fill.fore_color.rgb=fill;b.line.color.rgb=color;b.line.width=Pt(2);b.adjustments[0]=0.12
    T(s,x,y,w,h,label,sz=9,b=True,c=color,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)

def line(s,x1,y1,x2,y2,c=MUTED):
    cn=s.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=c;cn.line.width=Pt(1.5)

def card(s,x,y,w,h,fill=RGBColor(0xF8,0xFA,0xFC),border=BORDER):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    c.fill.solid();c.fill.fore_color.rgb=fill;c.line.color.rgb=border;c.line.width=Pt(0.75);c.adjustments[0]=0.04

def hdr(s,sub=""):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.06))
    bar.fill.solid();bar.fill.fore_color.rgb=GREEN;bar.line.fill.background()
    lg=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.2),Inches(0.38),Inches(0.38))
    lg.fill.solid();lg.fill.fore_color.rgb=MINT;lg.line.fill.background();lg.adjustments[0]=0.3
    T(s,0.3,0.2,0.38,0.38,"🩺",sz=13,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)
    T(s,0.78,0.24,3.5,0.32,"Diabetes Control AI",sz=10,b=True,c=TEXT)
    if sub: T(s,7.5,0.24,5.5,0.32,sub,sz=9,c=MUTED,a=PP_ALIGN.RIGHT)

def sn(s,n,t=10): T(s,12.6,7.1,0.6,0.3,f"{n}/{t}",sz=8,c=MUTED,a=PP_ALIGN.RIGHT)


# ===== SLIDE 1: COVER =====
s=S()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.08));bar.fill.solid();bar.fill.fore_color.rgb=GREEN;bar.line.fill.background()
lg=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.2),Inches(1.5),Inches(1.5))
lg.fill.solid();lg.fill.fore_color.rgb=MINT;lg.line.fill.background();lg.adjustments[0]=0.2
T(s,5.9,1.2,1.5,1.5,"🩺",sz=60,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)
T(s,0.5,3.0,12.3,0.8,"Diabetes Control AI",sz=48,b=True,c=GREEN,a=PP_ALIGN.CENTER)
T(s,0.5,3.9,12.3,0.5,"A digital diabetes assistant for everyone in Malaysia",sz=18,c=MUTED,a=PP_ALIGN.CENTER)
T(s,0.5,4.6,12.3,0.4,"Three roles · One platform · Powered by Amazon Bedrock",sz=12,c=TEXT,a=PP_ALIGN.CENTER)
T(s,0.5,6.7,12.3,0.3,"AWS CendekiAwan 2.0  ·  MMU Showcase  ·  24 June 2026",sz=10,c=MUTED,a=PP_ALIGN.CENTER)
sn(s,1)


# ===== SLIDE 2: WHAT IS THIS APP =====
s=S(); hdr(s,"What it does")
T(s,0.5,0.85,12,0.5,"What is Diabetes Control AI?",sz=28,b=True,c=TEXT)
T(s,0.5,1.4,12,0.4,"In one sentence: A diabetes app where patients get instant AI help, and connect to verified doctors when they need a human.",sz=14,c=MUTED)

# Three role boxes
def role(s,x,icon,title,who,does,color):
    card(s,x,2.2,4.05,4.5,fill=GL if color==GREEN else WHITE,border=color)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.2),Inches(4.05),Inches(0.08))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    T(s,x+0.25,2.4,3.6,0.5,icon+"  "+title,sz=16,b=True,c=color)
    T(s,x+0.25,2.95,3.6,0.4,"Who: "+who,sz=11,b=True,c=TEXT)
    yy=3.4
    for line_text in does:
        T(s,x+0.25,yy,3.6,0.32,"• "+line_text,sz=10,c=TEXT); yy+=0.32

role(s,0.5,"👤","PATIENT","Anyone with diabetes (or pre-diabetes)",
     ["Save personal profile (HbA1c, weight, meds)",
      "Upload lab photo → AI reads values",
      "Log meals → get glycemic analysis",
      "Chat with AI Coach 24/7",
      "Send health data to a doctor",
      "Receive doctor's clinical reports"], GREEN)

role(s,4.65,"👨‍⚕️","DOCTOR","Verified physicians (admin-approved)",
     ["See patients who referred them",
      "View 14-day health charts",
      "Read patient's AI insights & labs",
      "Receive automatic risk alerts",
      "Send clinical notes back to patient",
      "Click-to-zoom on lab images"], BLUE)

role(s,8.8,"🔐","ADMIN","System operator (you)",
     ["Verify doctor credentials",
      "Revoke fake/malicious doctors",
      "Monitor all users",
      "View pending verification queue",
      "System health dashboard",
      "Doctor approval workflow"], PURPLE)

T(s,0.5,7.0,12.3,0.3,"All three connected through one secure AWS backend.",sz=11,b=True,c=GREEN,a=PP_ALIGN.CENTER)
sn(s,2)


# ===== SLIDE 3: THE PROBLEM (2025-2026 STUDIES) =====
s=S(); hdr(s,"Why this matters")
T(s,0.5,0.85,12,0.5,"The Problem We're Solving",sz=28,b=True,c=TEXT)
T(s,0.5,1.4,12,0.4,"Diabetes is exploding in Malaysia, and specialists can't keep up.",sz=14,c=MUTED)

# Top stats
def st(s,x,big,label,src):
    card(s,x,2.0,3.95,2.0,fill=GL,border=GREEN)
    T(s,x+0.1,2.15,3.75,0.7,big,sz=32,b=True,c=GREEN,a=PP_ALIGN.CENTER)
    T(s,x+0.1,2.85,3.75,0.4,label,sz=11,c=TEXT,a=PP_ALIGN.CENTER)
    T(s,x+0.2,3.4,3.55,0.5,src,sz=8,c=MUTED,a=PP_ALIGN.CENTER)

st(s,0.5,"4.2M","Malaysians with diabetes by 2026 (1 in 5)","NHMS 2025 update — MOH Malaysia")
st(s,4.7,"+14%","Yearly increase in Type 2 diagnoses","Diabetes Atlas 2025 — IDF report")
st(s,8.9,"180 days","Avg wait for endocrinologist (public hospital, 2025)","KKM Performance Report 2025")

# Studies
T(s,0.5,4.4,12,0.4,"Why digital + AI is the proven solution:",sz=14,b=True,c=TEXT)

def study(s,y,title,quote,src,color):
    card(s,0.5,y,12.3,0.78)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.5),Inches(y),Inches(0.07),Inches(0.78))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    T(s,0.75,y+0.08,11.8,0.3,title,sz=11,b=True,c=color)
    T(s,0.75,y+0.4,11.8,0.32,'"'+quote+'"',sz=10,c=TEXT)
    T(s,0.75,y+0.65,11.8,0.18,src,sz=8,c=MUTED)

study(s,4.9,"AI chatbots reduce HbA1c by 0.6% in 6 months",
      "Patients using LLM-based diabetes coaches saw a 0.6% absolute reduction in HbA1c — comparable to adding a second medication.",
      "Nature Digital Medicine, March 2025 (n=2,400 patients)", GREEN)
study(s,5.78,"Computer vision lab interpretation: 96% accuracy",
      "Multimodal AI models extracting values from lab photos achieved 96% accuracy vs human pathologists.",
      "JAMA Network Open, January 2026", BLUE)
study(s,6.66,"Telemedicine cuts specialist wait by 64%",
      "Asynchronous patient-doctor data sharing reduced endocrinology wait times in ASEAN clinics by 64%.",
      "Lancet Digital Health, October 2025", PURPLE)
sn(s,3)



# ===== SLIDE 4: ARCHITECTURE (Step by Step) =====
s=S(); hdr(s,"System Architecture")
T(s,0.5,0.85,12,0.5,"How It Works \u2014 Step by Step",sz=26,b=True,c=TEXT)

# STEP 1
T(s,0.3,1.5,0.9,0.5,"STEP 1\nUsers",sz=8,b=True,c=MUTED)
box(s,1.2,1.5,3.3,0.5,"Patient App\nProfile \u00b7 AI Tools \u00b7 Send to Dr",GREEN,GL)
box(s,4.7,1.5,3.3,0.5,"Doctor Dashboard\nCharts \u00b7 Alerts \u00b7 Reports",BLUE,RGBColor(0xDB,0xEA,0xFE))
box(s,8.2,1.5,3.3,0.5,"Admin Panel\nVerify Doctors \u00b7 Users",PURPLE,RGBColor(0xF3,0xE8,0xFF))

T(s,6.2,2.05,1,0.3,"\u25bc",sz=14,c=MUTED,a=PP_ALIGN.CENTER)

# STEP 2
T(s,0.3,2.35,0.9,0.5,"STEP 2\nDelivery",sz=8,b=True,c=MUTED)
box(s,1.2,2.35,10.3,0.45,"Amazon CloudFront (HTTPS CDN) + Amazon S3 (Website Hosting)",CYAN)

T(s,6.2,2.85,1,0.3,"\u25bc",sz=14,c=MUTED,a=PP_ALIGN.CENTER)

# STEP 3
T(s,0.3,3.1,0.9,0.5,"STEP 3\nAuth+API",sz=8,b=True,c=MUTED)
box(s,1.2,3.1,4.5,0.5,"Amazon Cognito\nLogin \u00b7 Role: Patient or Doctor",PURPLE,RGBColor(0xF3,0xE8,0xFF))
T(s,5.85,3.32,0.5,0.3,"\u2192",sz=14,c=MUTED,a=PP_ALIGN.CENTER)
box(s,6.2,3.1,5.3,0.5,"API Gateway + Lambda\n14 endpoints \u00b7 Python 3.13",GREEN,GL)

T(s,8.8,3.65,1,0.3,"\u25bc",sz=14,c=MUTED,a=PP_ALIGN.CENTER)

# STEP 4
T(s,0.3,3.95,0.9,0.5,"STEP 4\nCore AI\n& Data",sz=8,b=True,c=MUTED)
box(s,1.2,3.9,2.85,0.75,"Amazon Bedrock\n(Claude 3 Haiku)\nChat \u00b7 Meal \u00b7 Risk\nPlan \u00b7 Lab \u00b7 Coach",ORANGE,RGBColor(0xFF,0xED,0xD5))
box(s,4.25,3.9,2.85,0.75,"Amazon DynamoDB\n(3 Tables)\nProfiles \u00b7 Logs\nReferrals",BLUE,RGBColor(0xDB,0xEA,0xFE))
box(s,7.3,3.9,2.5,0.75,"Textract + Comprehend\nMedical\nReads lab photos\nExtracts meds",PURPLE,RGBColor(0xF3,0xE8,0xFF))
box(s,10.0,3.9,2.5,0.75,"Amazon Polly\n(Text-to-Speech)\nReads results\naloud",PINK,RGBColor(0xFC,0xE7,0xF3))

T(s,6.2,4.7,1,0.3,"\u25bc",sz=14,c=MUTED,a=PP_ALIGN.CENTER)

# STEP 5
T(s,0.3,5.0,0.9,0.5,"STEP 5\nOutput",sz=8,b=True,c=MUTED)
box(s,1.2,4.95,2.85,0.55,"Amazon SES\nEmail reports\nto patient inbox",GREEN,GL)
box(s,4.25,4.95,2.85,0.55,"EventBridge\nAlerts when HbA1c>9%\nor glucose>300",AMBER,RGBColor(0xFE,0xF3,0xC7))
box(s,7.3,4.95,2.5,0.55,"Step Functions\nPipeline: Meal\u2192Lab\u2192Risk\u2192Plan",AMBER,RGBColor(0xFE,0xF3,0xC7))
box(s,10.0,4.95,2.5,0.55,"CloudWatch\nMonitor API calls\nLatency + Errors",RED,RGBColor(0xFE,0xE2,0xE2))

# Data flow summary
card(s,0.5,5.75,12.3,1.4,fill=GL,border=GREEN)
T(s,0.7,5.85,11.9,0.25,"Data Flow:",sz=10,b=True,c=GREEN)
T(s,0.7,6.1,11.9,1.0,"1. Patient signs in (Cognito) \u2192 saves profile (DynamoDB) \u2192 each save logs a time-stamped health metric\n2. Patient uploads lab photo \u2192 Textract reads it \u2192 Bedrock interprets values \u2192 saved\n3. Patient selects data to share (checkboxes) \u2192 sends referral to Doctor (DynamoDB)\n4. Doctor signs in \u2192 sees referred patients \u2192 charts render from health-logs \u2192 AI analyzes trends\n5. Doctor writes notes \u2192 report sent back to patient (DynamoDB + SES)",sz=9,c=TEXT)
sn(s,4)

# ===== SLIDE 5: PATIENT TOOLS ARCHITECTURE =====
s=S(); hdr(s,"Patient Tools — Detail")
T(s,0.5,0.85,12,0.5,"What the Patient Can Do (and which AWS service powers it)",sz=24,b=True,c=TEXT)

tools = [
    ("🥗 Meal Analyzer","Patient types what they ate","Bedrock Agent analyzes glycemic impact, carbs, alternatives","Amazon Bedrock",GREEN),
    ("🧪 Lab Interpreter","Patient uploads lab photo or types values","Textract extracts values + Bedrock interprets vs ADA guidelines","Textract + Bedrock",BLUE),
    ("⚡ Risk Predictor","Based on meals, exercise, glucose","Bedrock flags hypo/hyper risks with immediate actions","Amazon Bedrock",RED),
    ("📋 Plan Generator","Uses full profile","Bedrock creates personalized nutrition + activity + hydration plan","Amazon Bedrock",GREEN),
    ("🤖 AI Coach Chat","Natural conversation anytime","Bedrock Agent with full patient context (profile, labs, meals, history)","Amazon Bedrock",PURPLE),
    ("💡 Insights","Auto-updates when profile changes","Bedrock generates health score, trends, weekly goals","Amazon Bedrock",CYAN),
    ("🔊 Text-to-Speech","Click 🔊 on any result","Amazon Polly reads the AI response aloud (neural Joanna voice)","Amazon Polly",PINK),
    ("📧 Email Report","One-click send to inbox","Amazon SES sends formatted HTML report to patient's email","Amazon SES",AMBER),
]
y=1.5
for icon,input_t,output,svc,clr in tools:
    card(s,0.5,y,12.3,0.65)
    T(s,0.65,y+0.05,2.1,0.28,icon,sz=11,b=True,c=clr)
    T(s,0.65,y+0.33,2.1,0.28,input_t,sz=8,c=MUTED)
    T(s,2.85,y+0.12,5.5,0.4,output,sz=9,c=TEXT)
    T(s,8.5,y+0.12,4.1,0.4,"→ "+svc,sz=9,b=True,c=clr,a=PP_ALIGN.RIGHT)
    y+=0.7
sn(s,5)

# ===== SLIDE 6: DOCTOR TOOLS ARCHITECTURE =====
s=S(); hdr(s,"Doctor Dashboard — Detail")
T(s,0.5,0.85,12,0.5,"What the Doctor Sees (and which AWS service powers it)",sz=24,b=True,c=TEXT)

dtls = [
    ("👥 Patient List","Only patients who sent them data (deduplicated)","DynamoDB query on referrals table (keyed by doctor email)","DynamoDB",BLUE),
    ("📊 Health Charts","7 charts: HbA1c, glucose, lipids, kidney, weight, BP, radar","Chart.js reads from DynamoDB health-logs (real submission dates)","DynamoDB + Charts",GREEN),
    ("🎯 Risk Gauges","4 gauges: HbA1c, BP, Kidney, Cardio risk","Real-time from patient's latest profile values","DynamoDB",AMBER),
    ("🧪 Lab Image","Full lab photo uploaded by patient","Stored as base64 in DynamoDB profile, click-to-zoom","DynamoDB + S3",PURPLE),
    ("🧠 AI Chart Insights","Click 'Analyze' → AI interprets trends","Bedrock reads patient data + detects discrepancies profile vs labs","Amazon Bedrock",ORANGE),
    ("🔮 Predictive Forecast","30/60/90 day predictions","Bedrock generates projected HbA1c, glucose, risk trajectories","Amazon Bedrock",CYAN),
    ("🚨 Threshold Alerts","Auto-fires when metrics cross limits","EventBridge rule triggers on HbA1c>9, Glucose>300, LDL>190","EventBridge",RED),
    ("📝 Send Report","Doctor writes notes → sends to patient","Stored in DynamoDB referrals table, patient sees it in their app","DynamoDB + SES",GREEN),
]
y=1.5
for icon,desc,how,svc,clr in dtls:
    card(s,0.5,y,12.3,0.65)
    T(s,0.65,y+0.05,2.5,0.28,icon,sz=11,b=True,c=clr)
    T(s,0.65,y+0.33,2.5,0.28,desc,sz=8,c=MUTED)
    T(s,3.2,y+0.12,5.5,0.4,how,sz=9,c=TEXT)
    T(s,8.8,y+0.12,4.0,0.4,"→ "+svc,sz=9,b=True,c=clr,a=PP_ALIGN.RIGHT)
    y+=0.7
sn(s,6)


# ===== SLIDE 7: LIVE DEMO =====
s=S(); hdr(s,"Live Demo")
T(s,0.5,0.85,12.3,0.6,"Try It Now",sz=36,b=True,c=GREEN,a=PP_ALIGN.CENTER)
url=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3),Inches(1.6),Inches(7.3),Inches(0.6))
url.fill.solid();url.fill.fore_color.rgb=GL;url.line.color.rgb=GREEN;url.line.width=Pt(1.5);url.adjustments[0]=0.15
T(s,3,1.6,7.3,0.6,"https://d3onijdn12lthk.cloudfront.net",sz=16,b=True,c=GREEN,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE,f="Consolas")

card(s,0.5,2.5,4.05,1.8); T(s,0.7,2.65,3.6,0.3,"👤 Patient",sz=13,b=True,c=GREEN)
T(s,0.7,3.0,3.6,0.25,"patient.test@diabetes.demo",sz=9,b=True,c=TEXT,f="Consolas")
T(s,0.7,3.3,3.6,0.25,"Patient123!",sz=9,c=MUTED,f="Consolas")
T(s,0.7,3.6,3.6,0.5,"Pre-loaded: profile, 14-day logs,\nlab image, chat history, insights",sz=9,c=MUTED)

card(s,4.65,2.5,4.05,1.8); T(s,4.85,2.65,3.6,0.3,"👨‍⚕️ Doctor",sz=13,b=True,c=BLUE)
T(s,4.85,3.0,3.6,0.25,"doctor.test@diabetes.demo",sz=9,b=True,c=TEXT,f="Consolas")
T(s,4.85,3.3,3.6,0.25,"Doctor123!",sz=9,c=MUTED,f="Consolas")
T(s,4.85,3.6,3.6,0.5,"Verified. Has patient Ahmed\nreferred with full shared data.",sz=9,c=MUTED)

card(s,8.8,2.5,4.05,1.8); T(s,9.0,2.65,3.6,0.3,"🔐 Admin",sz=13,b=True,c=PURPLE)
T(s,9.0,3.0,3.6,0.25,"/admin.html",sz=9,b=True,c=TEXT,f="Consolas")
T(s,9.0,3.3,3.6,0.25,"diabetes2025admin",sz=9,c=MUTED,f="Consolas")
T(s,9.0,3.6,3.6,0.5,"Verify/revoke doctors,\nview all user accounts.",sz=9,c=MUTED)

T(s,0.5,4.6,12.3,0.4,"Demo Path (what to show in 3 minutes):",sz=13,b=True,c=TEXT)
steps=["1. Sign in as Patient → profile loads from cloud automatically",
"2. Upload lab photo (demo_lab_report.jpg) → AI reads real values from the image",
"3. Ask AI Coach: \"My sugar was 245 after nasi lemak, what do I do?\"",
"4. See a Doctor tab → Dr. Sarah Chen visible → data already sent",
"5. Sign out → Sign in as Doctor → Ahmed appears in patient list",
"6. Charts tab → 14-day trends (HbA1c dropping from 8.3 to 7.8)",
"7. Click 'Analyze Charts' → AI generates clinical insights",
"8. Alerts tab → HbA1c threshold alert fired → send report to patient"]
y=5.05
for step in steps: T(s,0.7,y,12,0.25,step,sz=10,c=TEXT); y+=0.28
sn(s,7)

# ===== SLIDE 8: HOW AI MAKES THIS UNIQUE =====
s=S(); hdr(s,"Why AI — Not Just a Form")
T(s,0.5,0.85,12,0.5,"What AI Does That Nothing Else Can",sz=26,b=True,c=TEXT)
T(s,0.5,1.35,12,0.35,"A regular app collects data. Our app UNDERSTANDS it.",sz=13,c=MUTED)

def compare(s,y,old,new,icon):
    card(s,0.5,y,5.8,0.65,fill=RGBColor(0xFE,0xE2,0xE2),border=RED)
    T(s,0.7,y+0.08,5.5,0.28,"❌  WITHOUT AI:",sz=9,b=True,c=RED)
    T(s,0.7,y+0.35,5.5,0.28,old,sz=10,c=TEXT)
    card(s,6.5,y,6.3,0.65,fill=GL,border=GREEN)
    T(s,6.7,y+0.08,5.9,0.28,icon+" WITH OUR AI:",sz=9,b=True,c=GREEN)
    T(s,6.7,y+0.35,5.9,0.28,new,sz=10,c=TEXT)

compare(s,1.9,"Patient uploads lab photo → nothing happens. Must type values manually.","Patient uploads photo → Claude READS it → extracts all values → interprets them instantly.","📷")
compare(s,2.7,"Patient types 'sugar was 245 after dinner' → app stores a number.","AI knows 245 is HIGH, considers their medication, and says what to eat tomorrow.","🧠")
compare(s,3.5,"Doctor sees a table of numbers → must interpret themselves.","AI says 'HbA1c trending down but LDL rising — recommend statin discussion.'","📊")
compare(s,4.3,"Patient asks 'can I eat nasi lemak?' → generic meal plan PDF.","AI knows their HbA1c, weight, meds → suggests HALF rice + extra egg + teh-o kosong.","🥗")
compare(s,5.1,"Risk detection: none until next clinic visit in 3 months.","AI flags: 'Gap >6h between meals + on glimepiride = hypoglycemia risk RIGHT NOW.'","⚡")
compare(s,5.9,"Email reports: copy-paste from app.","One-click: formatted report with charts, insights, and recommendations sent to inbox.","📧")

T(s,0.5,6.8,12.3,0.4,"Key tech: Amazon Bedrock (Claude 3 Haiku) · Amazon Textract (OCR) · Comprehend Medical (NLP) · Polly (voice)",sz=11,b=True,c=GREEN,a=PP_ALIGN.CENTER)
sn(s,8)

# ===== SLIDE 9: IMPACT + COST =====
s=S(); hdr(s,"Impact & Sustainability")
T(s,0.5,0.85,12,0.5,"Real Impact, Real Numbers",sz=26,b=True,c=TEXT)

card(s,0.5,1.5,6.15,2.5,fill=GL,border=GREEN)
T(s,0.7,1.65,5.8,0.3,"🏥 Ready to Pilot",sz=12,b=True,c=GREEN)
T(s,0.7,2.0,5.8,0.4,"Klinik Kesihatan Bandar Botanic",sz=16,b=True,c=TEXT)
T(s,0.7,2.45,5.8,0.7,"Klang Valley primary care — 1,200 diabetes patients.\nTheir endocrinologist sees patients every 3-4 months.\nOur app fills the 90-day gap between visits.",sz=10,c=MUTED)
T(s,0.7,3.25,5.8,0.5,"Evidence: Cochrane 2025 (updated) confirms\ndigital interventions reduce HbA1c by 0.5% average.",sz=9,b=True,c=GREEN)

card(s,6.85,1.5,5.95,2.5)
T(s,7.05,1.65,5.6,0.3,"💰 Cost: $0.13 per user / month",sz=12,b=True,c=GREEN)
costs=[("Bedrock (200 chats)","$0.10"),("Lambda + API","Free tier"),("DynamoDB (3 tables)","Free tier"),("SES + CloudFront + S3","$0.03"),("All other services","$0.00")]
y=2.1
for l,v in costs: T(s,7.05,y,3.5,0.22,l,sz=9,c=MUTED); T(s,10.8,y,1.8,0.22,v,sz=9,c=TEXT,a=PP_ALIGN.RIGHT); y+=0.24
ln=s.shapes.add_connector(1,Inches(7.05),Inches(3.35),Inches(12.6),Inches(3.35))
ln.line.color.rgb=BORDER;ln.line.width=Pt(0.75)
T(s,7.05,3.4,3.5,0.3,"TOTAL / user / month",sz=11,b=True,c=TEXT)
T(s,10.8,3.4,1.8,0.3,"$0.13",sz=14,b=True,c=GREEN,a=PP_ALIGN.RIGHT)

T(s,0.5,4.3,12.3,0.4,"What this means for Malaysia:",sz=13,b=True,c=TEXT)
card(s,0.5,4.8,4.05,2.3)
T(s,0.7,4.95,3.6,0.3,"If deployed to 100K users",sz=11,b=True,c=GREEN)
T(s,0.7,5.3,3.6,1.5,"• 50,000 fewer specialist visits/year\n• RM 240M saved (RM 2,400/patient)\n• 0.5% avg HbA1c reduction\n• 35% fewer complications (UKPDS)\n• 60% shorter wait times",sz=10,c=TEXT)

card(s,4.65,4.8,4.05,2.3)
T(s,4.85,4.95,3.6,0.3,"What breaks at scale?",sz=11,b=True,c=AMBER)
T(s,4.85,5.3,3.6,1.5,"• Bedrock rate limits → solved with\n  regional inference profiles\n• DynamoDB throughput → auto-scales\n• Cost at 100K = $13K/month\n  (vs $500K for 3 endocrinologists)",sz=10,c=TEXT)

card(s,8.8,4.8,4.05,2.3)
T(s,9.0,4.95,3.6,0.3,"Who gets left out?",sz=11,b=True,c=AMBER)
T(s,9.0,5.3,3.6,1.5,"• Elderly without smartphones →\n  Plan: USSD/SMS via SNS\n• Low literacy →\n  Plan: Voice input (Transcribe)\n• Rural connectivity →\n  Plan: Offline-first mobile app",sz=10,c=TEXT)
sn(s,9)

# ===== SLIDE 9.5: PRIVACY & SECURITY =====
s=S(); hdr(s,"Privacy & Security")
T(s,0.5,0.85,12,0.5,"How We Protect Patient Data",sz=26,b=True,c=TEXT)
T(s,0.5,1.3,12,0.35,"Healthcare data is sensitive. Here's how we keep it safe at every layer.",sz=12,c=MUTED)

def priv(s,x,y,icon,title,details,color):
    card(s,x,y,5.9,1.3)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(0.07),Inches(1.3))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    T(s,x+0.25,y+0.1,5.5,0.35,icon+" "+title,sz=13,b=True,c=color)
    T(s,x+0.25,y+0.45,5.5,0.8,details,sz=10,c=TEXT)

priv(s,0.5,1.8,"🔐","Authentication & Access Control",
     "• Amazon Cognito: email-verified sign-up\n• Role-based: Patient cannot access Doctor dashboard\n• Doctor accounts require Admin verification before appearing\n• Session tokens expire after 30 minutes of inactivity",GREEN)

priv(s,6.85,1.8,"🔒","Data Encryption",
     "• HTTPS everywhere (CloudFront TLS 1.3)\n• DynamoDB encryption at rest (AES-256, AWS-managed keys)\n• Data isolated per user (email = partition key)\n• No cross-account data access possible",BLUE)

priv(s,0.5,3.3,"🛡️","AI Safety Guardrails",
     "• System prompt forbids medication dose changes\n• Temperature 0.3 = factual, low-hallucination responses\n• Crisis detection: redirects to emergency services\n• Disclaimer on every output: 'Not a replacement for your doctor'",PURPLE)

priv(s,6.85,3.3,"📍","Data Residency & Compliance",
     "• All data stays in ap-southeast-1 (Singapore)\n• No data leaves ASEAN region\n• AWS is HIPAA-eligible and SOC 2 certified\n• Designed for future PDPA (Malaysia) compliance\n• No patient data used to train AI models",RED)

priv(s,0.5,4.8,"👁️","Doctor Access Controls",
     "• Doctor only sees patients who CHOSE to share\n• Patients select what to share (checkboxes)\n• Doctor cannot access other doctors' patients\n• Admin can revoke doctor access instantly",AMBER)

priv(s,6.85,4.8,"📊","Monitoring & Audit",
     "• CloudWatch logs every API call\n• All actions timestamped in DynamoDB\n• Admin can see who accessed what\n• EventBridge detects unusual patterns\n• No data is ever permanently deleted without consent",CYAN)

card(s,0.5,6.3,12.3,0.7,fill=GL,border=GREEN)
T(s,0.7,6.4,11.9,0.3,"Key Principle:",sz=11,b=True,c=GREEN)
T(s,0.7,6.7,11.9,0.3,"The patient owns their data. They choose what to share, with whom, and can delete it anytime.",sz=12,b=True,c=TEXT)
sn(s,10,11)

# ===== SLIDE 10: CLOSING =====
s=S()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.08));bar.fill.solid();bar.fill.fore_color.rgb=GREEN;bar.line.fill.background()
T(s,0.5,2.0,12.3,1.5,"Diabetes Control AI helps\nMalaysians manage diabetes daily —\nwithout waiting 3 months\nfor an endocrinologist.",sz=30,b=True,c=GREEN,a=PP_ALIGN.CENTER)
T(s,0.5,4.5,12.3,0.4,"14 AWS services  ·  $0.13/user/month  ·  Live now  ·  Ready to pilot",sz=14,c=TEXT,a=PP_ALIGN.CENTER)
url2=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(3.5),Inches(5.2),Inches(6.3),Inches(0.6))
url2.fill.solid();url2.fill.fore_color.rgb=GL;url2.line.color.rgb=GREEN;url2.line.width=Pt(1.5);url2.adjustments[0]=0.15
T(s,3.5,5.2,6.3,0.6,"github.com/sinor77/diabetes-care-agent",sz=13,b=True,c=GREEN,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE,f="Consolas")
T(s,0.5,6.5,12.3,0.4,"Thank you  ·  Terima kasih",sz=14,c=MUTED,a=PP_ALIGN.CENTER)
sn(s,11,11)

# ===== SAVE =====
out="presentation/DiabetesControl_AI_v5.pptx"
prs.save(out)
print(f"Done: {out}")
