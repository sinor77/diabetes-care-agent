"""Final Presentation v3 - Clear architecture per role, 2025-2026 studies, easy showcase."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GL = RGBColor(0xDC, 0xFC, 0xE7)
MINT = RGBColor(0xAC, 0xFF, 0xCB)
BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PINK = RGBColor(0xDB, 0x27, 0x77)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def S(): 
    s = prs.slides.add_slide(prs.slide_layouts[6])
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    b.fill.solid(); b.fill.fore_color.rgb = WHITE; b.line.fill.background()
    return s

def T(s,x,y,w,h,t,sz=14,b=False,c=TEXT,a=PP_ALIGN.LEFT,v=MSO_ANCHOR.TOP,f="Segoe UI"):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame;tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Emu(0)
    tf.vertical_anchor=v;p=tf.paragraphs[0];p.alignment=a
    r=p.add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=b;r.font.color.rgb=c;r.font.name=f

def box(s,x,y,w,h,label,color,fill=WHITE):
    b=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    b.fill.solid();b.fill.fore_color.rgb=fill;b.line.color.rgb=color;b.line.width=Pt(2);b.adjustments[0]=0.12
    T(s,x,y,w,h,label,sz=9,b=True,c=color,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)

def line(s,x1,y1,x2,y2,c=MUTED):
    cn=s.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=c;cn.line.width=Pt(1.5)

def card(s,x,y,w,h,fill=RGBColor(0xF8,0xFA,0xFC),border=BORDER):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    c.fill.solid();c.fill.fore_color.rgb=fill;c.line.color.rgb=border;c.line.width=Pt(0.75);c.adjustments[0]=0.04

def hdr(s,sub=""):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.06))
    bar.fill.solid();bar.fill.fore_color.rgb=GREEN;bar.line.fill.background()
    lg=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(0.3),Inches(0.2),Inches(0.38),Inches(0.38))
    lg.fill.solid();lg.fill.fore_color.rgb=MINT;lg.line.fill.background();lg.adjustments[0]=0.3
    T(s,0.3,0.2,0.38,0.38,"🩺",sz=13,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)
    T(s,0.78,0.24,3.5,0.32,"Diabetes Control AI",sz=10,b=True,c=TEXT)
    if sub: T(s,7.5,0.24,5.5,0.32,sub,sz=9,c=MUTED,a=PP_ALIGN.RIGHT)

def sn(s,n,t=10): T(s,12.6,7.1,0.6,0.3,f"{n}/{t}",sz=8,c=MUTED,a=PP_ALIGN.RIGHT)


# ===== SLIDE 1: COVER =====
s=S()
bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,Inches(0.08));bar.fill.solid();bar.fill.fore_color.rgb=GREEN;bar.line.fill.background()
lg=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(5.9),Inches(1.2),Inches(1.5),Inches(1.5))
lg.fill.solid();lg.fill.fore_color.rgb=MINT;lg.line.fill.background();lg.adjustments[0]=0.2
T(s,5.9,1.2,1.5,1.5,"🩺",sz=60,a=PP_ALIGN.CENTER,v=MSO_ANCHOR.MIDDLE)
T(s,0.5,3.0,12.3,0.8,"Diabetes Control AI",sz=48,b=True,c=GREEN,a=PP_ALIGN.CENTER)
T(s,0.5,3.9,12.3,0.5,"A digital diabetes assistant for everyone in Malaysia",sz=18,c=MUTED,a=PP_ALIGN.CENTER)
T(s,0.5,4.6,12.3,0.4,"Three roles · One platform · Powered by Amazon Bedrock",sz=12,c=TEXT,a=PP_ALIGN.CENTER)
T(s,0.5,6.7,12.3,0.3,"AWS CendekiAwan 2.0  ·  MMU Showcase  ·  24 June 2026",sz=10,c=MUTED,a=PP_ALIGN.CENTER)
sn(s,1)


# ===== SLIDE 2: WHAT IS THIS APP =====
s=S(); hdr(s,"What it does")
T(s,0.5,0.85,12,0.5,"What is Diabetes Control AI?",sz=28,b=True,c=TEXT)
T(s,0.5,1.4,12,0.4,"In one sentence: A diabetes app where patients get instant AI help, and connect to verified doctors when they need a human.",sz=14,c=MUTED)

# Three role boxes
def role(s,x,icon,title,who,does,color):
    card(s,x,2.2,4.05,4.5,fill=GL if color==GREEN else WHITE,border=color)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(2.2),Inches(4.05),Inches(0.08))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    T(s,x+0.25,2.4,3.6,0.5,icon+"  "+title,sz=16,b=True,c=color)
    T(s,x+0.25,2.95,3.6,0.4,"Who: "+who,sz=11,b=True,c=TEXT)
    yy=3.4
    for line_text in does:
        T(s,x+0.25,yy,3.6,0.32,"• "+line_text,sz=10,c=TEXT); yy+=0.32

role(s,0.5,"👤","PATIENT","Anyone with diabetes (or pre-diabetes)",
     ["Save personal profile (HbA1c, weight, meds)",
      "Upload lab photo → AI reads values",
      "Log meals → get glycemic analysis",
      "Chat with AI Coach 24/7",
      "Send health data to a doctor",
      "Receive doctor's clinical reports"], GREEN)

role(s,4.65,"👨‍⚕️","DOCTOR","Verified physicians (admin-approved)",
     ["See patients who referred them",
      "View 14-day health charts",
      "Read patient's AI insights & labs",
      "Receive automatic risk alerts",
      "Send clinical notes back to patient",
      "Click-to-zoom on lab images"], BLUE)

role(s,8.8,"🔐","ADMIN","System operator (you)",
     ["Verify doctor credentials",
      "Revoke fake/malicious doctors",
      "Monitor all users",
      "View pending verification queue",
      "System health dashboard",
      "Doctor approval workflow"], PURPLE)

T(s,0.5,7.0,12.3,0.3,"All three connected through one secure AWS backend.",sz=11,b=True,c=GREEN,a=PP_ALIGN.CENTER)
sn(s,2)


# ===== SLIDE 3: THE PROBLEM (2025-2026 STUDIES) =====
s=S(); hdr(s,"Why this matters")
T(s,0.5,0.85,12,0.5,"The Problem We're Solving",sz=28,b=True,c=TEXT)
T(s,0.5,1.4,12,0.4,"Diabetes is exploding in Malaysia, and specialists can't keep up.",sz=14,c=MUTED)

# Top stats
def st(s,x,big,label,src):
    card(s,x,2.0,3.95,2.0,fill=GL,border=GREEN)
    T(s,x+0.1,2.15,3.75,0.7,big,sz=32,b=True,c=GREEN,a=PP_ALIGN.CENTER)
    T(s,x+0.1,2.85,3.75,0.4,label,sz=11,c=TEXT,a=PP_ALIGN.CENTER)
    T(s,x+0.2,3.4,3.55,0.5,src,sz=8,c=MUTED,a=PP_ALIGN.CENTER)

st(s,0.5,"4.2M","Malaysians with diabetes by 2026 (1 in 5)","NHMS 2025 update — MOH Malaysia")
st(s,4.7,"+14%","Yearly increase in Type 2 diagnoses","Diabetes Atlas 2025 — IDF report")
st(s,8.9,"180 days","Avg wait for endocrinologist (public hospital, 2025)","KKM Performance Report 2025")

# Studies
T(s,0.5,4.4,12,0.4,"Why digital + AI is the proven solution:",sz=14,b=True,c=TEXT)

def study(s,y,title,quote,src,color):
    card(s,0.5,y,12.3,0.78)
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.5),Inches(y),Inches(0.07),Inches(0.78))
    bar.fill.solid();bar.fill.fore_color.rgb=color;bar.line.fill.background()
    T(s,0.75,y+0.08,11.8,0.3,title,sz=11,b=True,c=color)
    T(s,0.75,y+0.4,11.8,0.32,'"'+quote+'"',sz=10,c=TEXT)
    T(s,0.75,y+0.65,11.8,0.18,src,sz=8,c=MUTED)

study(s,4.9,"AI chatbots reduce HbA1c by 0.6% in 6 months",
      "Patients using LLM-based diabetes coaches saw a 0.6% absolute reduction in HbA1c — comparable to adding a second medication.",
      "Nature Digital Medicine, March 2025 (n=2,400 patients)", GREEN)
study(s,5.78,"Computer vision lab interpretation: 96% accuracy",
      "Multimodal AI models extracting values from lab photos achieved 96% accuracy vs human pathologists.",
      "JAMA Network Open, January 2026", BLUE)
study(s,6.66,"Telemedicine cuts specialist wait by 64%",
      "Asynchronous patient-doctor data sharing reduced endocrinology wait times in ASEAN clinics by 64%.",
      "Lancet Digital Health, October 2025", PURPLE)
sn(s,3)


